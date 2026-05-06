"""PPTX extractor."""

from __future__ import annotations

from pathlib import Path

from dragndoc.config import get_settings
from dragndoc.extractors._caps import CapConfig, select_pages
from dragndoc.extractors._meta import collect
from dragndoc.extractors.base import CorruptDocumentError, ExtractedDoc, Section


# python-pptx CoreProperties attributes — same OOXML core schema as docx.
_CORE_ATTRS = (
    "title", "author", "subject", "keywords", "category", "comments",
    "last_modified_by", "revision", "version", "created", "modified",
    "last_printed", "content_status", "identifier", "language",
)


def extract(path: Path) -> ExtractedDoc:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise CorruptDocumentError("python-pptx is not installed") from exc
    try:
        pres = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        raise CorruptDocumentError(f"pptx failed for {path}: {exc}") from exc

    cfg = CapConfig.from_settings(get_settings())

    def _iter_slides():
        for slide in pres.slides:
            chunks: list[str] = []
            for shape in slide.shapes:
                tf = getattr(shape, "text_frame", None) if shape.has_text_frame else None
                if tf is not None:
                    for para in tf.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            chunks.append(line)
            yield "\n".join(chunks)

    kept = select_pages(_iter_slides(), cfg)
    sections = [
        Section(label=f"Slide {i + 1}", text=text, index=i)
        for i, text in enumerate(kept)
    ]

    raw: dict = {}
    try:
        cp = pres.core_properties
        for attr in _CORE_ATTRS:
            raw[attr] = getattr(cp, attr, None)
    except Exception:
        pass

    return ExtractedDoc(
        path=path,
        sections=sections,
        total_sections=len(pres.slides),
        format="pptx",
        extracted_metadata=collect(raw, prefix="core_"),
    )
