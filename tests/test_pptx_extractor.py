"""Tests for sectioned PPTX extraction."""

from __future__ import annotations

import pytest


def test_pptx_slides_become_sections(tmp_path):
    pptx_lib = pytest.importorskip("pptx")
    from pptx.util import Inches
    from dragndoc.extractors import pptx as pptx_ext

    path = tmp_path / "deck.pptx"
    pres = pptx_lib.Presentation()
    blank = pres.slide_layouts[6]
    for i in range(4):
        slide = pres.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text = f"Slide body {i + 1}"
    pres.save(path)

    doc = pptx_ext.extract(path)
    assert len(doc.sections) == 4
    assert doc.total_sections == 4
    assert [section.label for section in doc.sections] == [f"Slide {i}" for i in range(1, 5)]
    assert doc.sections[2].text == "Slide body 3"
