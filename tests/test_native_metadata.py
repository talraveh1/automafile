"""Tests for native metadata writers and mtime preservation."""

from __future__ import annotations

import os
import time
from pathlib import Path

import pytest

from automafile.metadata import native
from automafile.metadata.mtime import preserve_times, restore, snapshot


def _make_blank_pdf(path: Path) -> None:
    import pikepdf
    pdf = pikepdf.new()
    pdf.add_blank_page(page_size=(612, 792))
    pdf.save(path)


def _make_blank_docx(path: Path) -> None:
    from docx import Document
    doc = Document()
    doc.add_paragraph("hello")
    doc.save(str(path))


def _make_blank_xlsx(path: Path) -> None:
    from openpyxl import Workbook
    wb = Workbook()
    wb.active["A1"] = "hi"
    wb.save(str(path))


def _make_blank_pptx(path: Path) -> None:
    from pptx import Presentation
    pres = Presentation()
    pres.slides.add_slide(pres.slide_layouts[5])
    pres.save(str(path))


def _make_jpeg(path: Path) -> None:
    from PIL import Image
    Image.new("RGB", (16, 16), color=(255, 0, 0)).save(path, format="JPEG", quality=80)


def test_supports_known_formats(tmp_path):
    for ext in (".pdf", ".docx", ".xlsx", ".pptx", ".jpg", ".png"):
        assert native.supports(tmp_path / f"x{ext}")
    assert not native.supports(tmp_path / "x.txt")


def test_pdf_native_write_and_mtime_preserved(tmp_path):
    p = tmp_path / "doc.pdf"
    _make_blank_pdf(p)
    snap = snapshot(p)
    time.sleep(0.05)
    native.write(p, {
        "title": "Test Title",
        "summary": "Test summary written by automafile.",
        "tags": ["alpha", "beta"],
        "category": "Personal",
        "correspondent": "Acme",
        "date": "2026-04-01",
        "confidence": "high",
    })
    assert snapshot(p)[1] == snap[1]

    import pikepdf
    with pikepdf.open(p) as pdf:
        with pdf.open_metadata() as xmp:
            assert "Test Title" in xmp.get("dc:title", "")
            assert "Test summary" in xmp.get("dc:description", "")


def test_docx_native_write(tmp_path):
    p = tmp_path / "doc.docx"
    _make_blank_docx(p)
    snap = snapshot(p)
    time.sleep(0.05)
    native.write(p, {
        "title": "Hello",
        "summary": "A summary",
        "tags": ["one", "two"],
        "category": "Research",
    })
    assert snapshot(p)[1] == snap[1]
    from docx import Document
    cp = Document(str(p)).core_properties
    assert cp.title == "Hello"
    assert cp.subject == "A summary"
    assert cp.keywords == "one, two"
    assert cp.category == "Research"


def test_xlsx_native_write(tmp_path):
    p = tmp_path / "sheet.xlsx"
    _make_blank_xlsx(p)
    snap = snapshot(p)
    time.sleep(0.05)
    native.write(p, {
        "title": "T",
        "summary": "S",
        "tags": ["a"],
        "category": "Financial",
    })
    assert snapshot(p)[1] == snap[1]
    from openpyxl import load_workbook
    cp = load_workbook(str(p)).properties
    assert cp.title == "T"
    assert cp.category == "Financial"


def test_pptx_native_write(tmp_path):
    p = tmp_path / "deck.pptx"
    _make_blank_pptx(p)
    snap = snapshot(p)
    time.sleep(0.05)
    native.write(p, {
        "title": "Deck",
        "summary": "S",
        "tags": ["x"],
        "category": "Teaching",
    })
    assert snapshot(p)[1] == snap[1]


def test_jpeg_native_write(tmp_path):
    p = tmp_path / "img.jpg"
    _make_jpeg(p)
    snap = snapshot(p)
    time.sleep(0.05)
    native.write(p, {
        "title": "Photo",
        "summary": "A photo summary",
        "tags": ["pic"],
        "category": "Media",
    })
    assert snapshot(p)[1] == snap[1]


def test_preserve_times_context_manager(tmp_path):
    p = tmp_path / "x.txt"
    p.write_text("hi", encoding="utf-8")
    snap = snapshot(p)
    time.sleep(0.05)
    with preserve_times(p):
        p.write_text("changed", encoding="utf-8")
    assert snapshot(p)[1] == snap[1]
