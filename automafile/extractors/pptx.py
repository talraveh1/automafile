"""PPTX extractor."""

from __future__ import annotations

from pathlib import Path

from automafile.extractors._meta import collect
from automafile.extractors.base import CorruptDocumentError, ExtractedDoc


# python-pptx CoreProperties attributes — same OOXML core schema as docx.
_CORE_ATTRS = (
    "title", "author", "subject", "keywords", "category", "comments",
    "last_modified_by", "revision", "version", "created", "modified",
    "last_printed", "content_status", "identifier", "language",
)


def extract(path: Path) -> ExtractedDoc:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise CorruptDocumentError("python-pptx is not installed") from exc
    try:
        pres = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        raise CorruptDocumentError(f"pptx failed for {path}: {exc}") from exc

    chunks: list[str] = []
    for slide_idx, slide in enumerate(pres.slides, start=1):
        chunks.append(f"# Slide {slide_idx}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        chunks.append(line)

    raw: dict = {}
    try:
        cp = pres.core_properties
        for attr in _CORE_ATTRS:
            raw[attr] = getattr(cp, attr, None)
    except Exception:
        pass

    return ExtractedDoc(
        path=path,
        text="\n".join(chunks),
        format="pptx",
        extracted_metadata=collect(raw, prefix="core_"),
    )
